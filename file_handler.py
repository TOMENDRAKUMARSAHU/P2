
import os
import json
import pandas as pd
from docx import Document
from pptx import Presentation
import PyPDF2
from pathlib import Path

def extract_text_from_file(file_path):
    ext = Path(file_path).suffix.lower()

    if ext == ".txt":
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()

    elif ext == ".json":
        with open(file_path, "r", encoding="utf-8") as f:
            try:
                data = json.load(f)
                return json.dumps(data, indent=2)
            except Exception as e:
                return f"Error reading JSON: {e}"

    elif ext == ".csv":
        try:
            df = pd.read_csv(file_path)
            return df.to_string(index=False)
        except Exception as e:
            return f"Error reading CSV: {e}"

    elif ext == ".docx":
        try:
            doc = Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])
        except Exception as e:
            return f"Error reading DOCX: {e}"

    elif ext == ".pptx":
        try:
            prs = Presentation(file_path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)
        except Exception as e:
            return f"Error reading PPTX: {e}"

    elif ext == ".pdf":
        try:
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                return "\n".join([page.extract_text() or "" for page in reader.pages])
        except Exception as e:
            return f"Error reading PDF: {e}"

    else:
        return "Unsupported file type: " + ext

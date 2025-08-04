import os
import json
import pandas as pd
import zipfile
import csv
import docx
import pptx
import chardet
import PyPDF2
import textract
import re
import pyarrow.parquet as pq
import requests
from bs4 import BeautifulSoup


def parse_any_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == ".txt":
            with open(file_path, "rb") as f:
                encoding = chardet.detect(f.read())['encoding']
            with open(file_path, "r", encoding=encoding) as f:
                return f.read()

        elif ext == ".pdf":
            try:
                import pdfplumber
                with pdfplumber.open(file_path) as pdf:
                    return "\n".join(page.extract_text() or '' for page in pdf.pages)
            except:
                with open(file_path, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    return "\n".join(page.extract_text() or '' for page in reader.pages)

        elif ext == ".docx":
            doc = docx.Document(file_path)
            return "\n".join([p.text for p in doc.paragraphs])

        elif ext == ".pptx":
            ppt = pptx.Presentation(file_path)
            text = ""
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text

        elif ext == ".json":
            with open(file_path, "r", encoding="utf-8") as f:
                return json.dumps(json.load(f), indent=2)

        elif ext == ".csv":
            df = pd.read_csv(file_path)
            return df.to_string(index=False)

        elif ext in [".xls", ".xlsx"]:
            df = pd.read_excel(file_path)
            return df.to_string(index=False)

        elif ext == ".parquet":
            df = pd.read_parquet(file_path)
            return df.to_string(index=False)

        elif ext == ".zip":
            text = ""
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                for name in zip_ref.namelist():
                    with zip_ref.open(name) as f:
                        try:
                            content = f.read()
                            detected = chardet.detect(content)
                            text += content.decode(detected["encoding"] or "utf-8") + "\n"
                        except:
                            continue
            return text

        else:
            return textract.process(file_path).decode("utf-8")

    except Exception as e:
        return f"Error parsing file: {e}"


def extract_url_and_questions(text):
    urls = re.findall(r'(https?://\S+)', text)
    questions = re.findall(r'(?:Q[:\-\s]?|\n)([^\n?]{3,}\?)', text)
    return (urls[0] if urls else ""), questions


def convert_s3_to_https(s3_url):
    if not s3_url.startswith("s3://"):
        return s3_url
    s3_url = s3_url.replace("s3://", "")
    parts = s3_url.split("/", 1)
    bucket = parts[0]
    path = parts[1] if len(parts) > 1 else ""
    return f"https://{bucket}.s3.amazonaws.com/{path}"


def try_fetch_parquet_from_url(url):
    try:
        response = requests.get(url, stream=True)
        if response.status_code == 200:
            with open("temp.parquet", "wb") as f:
                f.write(response.content)
            df = pd.read_parquet("temp.parquet")
            os.remove("temp.parquet")
            return df.to_string(index=False)
        else:
            return f"Failed to download parquet: HTTP {response.status_code}"
    except Exception as e:
        return f"Error fetching parquet: {e}"
